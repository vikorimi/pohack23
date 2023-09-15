import os
import re
from googletrans import Translator
from pptx import Presentation
from pptx.util import Inches

text = """Q: where can we find the vlcentral taxonomy sheet?
A: Currently, there is no specific Taxonomy Sheet, Partner will need to select the right or related Taxonomy in My Cases and Support each time. We do not have a Taxonomy worksheet for VL Central and the guidance to navigate the resource is within the My Cases and Support materials. The VL Central case creation process enables the submitters  to submit their case based on their specific requirements The taxonomy is assigned in the back end based on the question that is being raised.
"""

languages = {
    'es': 'Spanish',
    'de': 'German',
    'ko': 'Korean',
    'pt': 'Portuguese',
    'zh-CN': 'Chinese'
}



def translate_text(text, target_language):
    translator = Translator()
    translated = translator.translate(text, dest=target_language)
    return translated.text

def extract_questions_and_answers(text):
    qa_pairs = re.findall(r'Q: (.*?)(?:\nA: |$)', text, re.DOTALL)
    questions_and_answers = []
    
    for question in qa_pairs:
        answer = re.search(fr'A: (.*?)(?=\nQ: |$)', text).group(1).strip()
        questions_and_answers.append((question.strip(), answer.strip()))
    
    return questions_and_answers

def create_powerpoint(qa_pairs, output_path, target_language):
    prs = Presentation()
    for i, (question, answer) in enumerate(qa_pairs, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout
        
        # Add the translated question
        translated_question = translate_text(question, target_language)
        question_text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        question_frame = question_text.text_frame
        question_frame.text = f"{languages[target_language]} Question {i}:"
        p = question_frame.add_paragraph()
        p.text = translated_question
        
        # Add the translated answer
        translated_answer = translate_text(answer, target_language)
        answer_text = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        answer_frame = answer_text.text_frame
        answer_frame.text = f"{languages[target_language]} Answer {i}:"
        p = answer_frame.add_paragraph()
        p.text = translated_answer
    
    prs.save(output_path)

def main():
    output_folder = "lsptranslated_presentations"
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    languages_to_translate = ['es', 'de', 'ko', 'pt', 'zh-CN']
    
    for lang_code in languages_to_translate:
        target_language = lang_code
        lang_name = languages[lang_code]
        output_path = os.path.join(output_folder, f"translated_presentation_{lang_name}.pptx")
        qa_pairs = extract_questions_and_answers(text)
        
        if not qa_pairs:
            print(f"No questions and answers found for {lang_name}.")
            continue
        
        create_powerpoint(qa_pairs, output_path, target_language)
        print(f"Presentation saved as {output_path} ({lang_name})")

if __name__ == "__main__":
    main()
