#!/usr/bin/env python
# coding: utf-8

# In[1]:


import pandas as pd

# Replace 'your_file.txt' with the path to your text file
file_path = './QnA.txt'

# Use read_csv to read the text file, specifying the delimiter (if it's not comma)
# For example, if it's a tab-separated file (TSV), you can use sep='\t'
# If the file has a header row, you can specify header=0 to use the first row as column names
# If there is no header, you can remove the header parameter or set it to None
df = pd.read_table(file_path, sep='\t', header=0)


# In[2]:


df[['Name', 'Question']] = df['header'].str.split(':', 1, expand=True)


# In[3]:


df['Answer'] = df['Question'].shift(-1)
df['WhoAnswer'] = df['Name'].shift(-1)


# In[4]:


df_1 = df


# In[5]:


df_1


# In[6]:


df_2=df_1[df_1.Name=="Ram"]


# In[7]:


df_2=df_2[["Name","Question","Answer","WhoAnswer"]]


# In[8]:


import pandas as pd

# Replace 'your_file.txt' with the path to your text file
file_path = './Transcript.txt'

# Use read_csv to read the text file, specifying the delimiter (if it's not comma)
# For example, if it's a tab-separated file (TSV), you can use sep='\t'
# If the file has a header row, you can specify header=0 to use the first row as column names
# If there is no header, you can remove the header parameter or set it to None
dfTranscript = pd.read_table(file_path, sep='\t', header=0)
dfTranscript[['Name', 'Question']] = dfTranscript['header'].str.split(':', 1, expand=True)
dfTranscript['Answer'] = dfTranscript['Question'].shift(-1)
dfTranscript['WhoAnswer'] = dfTranscript['Name'].shift(-1)
dfTranscript_1 = dfTranscript
dfTranscript_2=dfTranscript_1[dfTranscript_1.Name=="Prasad"]
dfTranscript_2=dfTranscript_2[["Name","Question","Answer","WhoAnswer"]]


# In[9]:


dfTranscript_2


# In[85]:





# In[10]:


df_all_data


# In[89]:





# In[14]:


import pandas as pd

# Specify the path to the text file you want to read
file_path = './Chat.txt'

# Initialize an empty list to store lines from the text file
lines = []

# Open the text file and read it line by line
with open(file_path, 'r') as file:
    for line in file:
        lines.append([line.strip()])  # Append each line as a list containing one element

# Create a Pandas DataFrame from the list of lines
df = pd.DataFrame(lines, columns=['Line'])
dfChat = df[df.Line != "header"]
dfChat
dfChat[['Name', 'Question']] = dfChat['Line'].str.split(':', 1, expand=True)
dfChat['Answer'] = dfChat['Question'].shift(-1)
dfChat['WhoAnswer'] = dfChat['Name'].shift(-1)
dfChat_1 = dfChat
dfChat_1
dfChat_2=dfChat_1[dfChat_1.Name=="CC"]
dfChat_2=dfChat_2[["Name","Question","Answer","WhoAnswer"]]


# In[15]:


df_all_data = pd.concat([df_2,dfTranscript_2,dfChat_2],axis =0)
df_all_data.rename(columns={'Question': 'Question'}, inplace=True)


# In[16]:


df_all_data


# In[127]:


df_all_data = pd.concat([df_2,dfTranscript_2,dfChat_2],axis =0)


excel_file_path = './output.xlsx'

# Write the DataFrame to an Excel file
df_all_data.to_excel(excel_file_path, index=False)  # Set index=False to exclude the index column


# In[ ]:





# In[124]:





# In[125]:


dfChat_2


# In[99]:





# In[100]:


df


# In[ ]:





# In[ ]:





# In[ ]:





# In[133]:


from pptx import Presentation


# In[135]:


from pptx import Presentation

# Create a PowerPoint presentation
prs = Presentation()

# Add a slide with Title Slide layout
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)

# Add a title and subtitle to the slide
title = slide.shapes.title
title.text = "My PowerPoint Slide"

subtitle = slide.placeholders[1]  # Subtitle placeholder
subtitle.text = "This is some text on the slide."

# Save the PowerPoint presentation to a file
prs.save("my_presentation.pptx")


# In[138]:


df_all_data


# In[136]:


# Specify the path to the existing PowerPoint presentation
pptx_file_path = './June FY23 - LSP Fundamentals Q&A (English).pptx'

# Load the PowerPoint presentation
prs = Presentation(pptx_file_path)

# Iterate through slides and extract information
for slide_number, slide in enumerate(prs.slides, start=1):
    print(f"Slide {slide_number}:")
    
    # Extract text from each shape on the slide
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            print(shape.text)
    
    print()


# In[150]:


import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
pptx_file_path = './June FY23 - LSP Fundamentals Q&A (English).pptx'

# Load the PowerPoint presentation
prs = Presentation(pptx_file_path)

# Add a slide
slide_layouts = prs.slide_master.slide_layouts
    
slide_layout = prs.slide_layouts[9]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)

# Add a title
title = slide.shapes.title
title.text = "DataFrame Table Example"

# Specify the position and dimensions of the table on the slide
left = Inches(1)
top = Inches(1)
width = Inches(10)
height = Inches(2)

# Add the DataFrame as a table to the slide
table = slide.shapes.add_table(df_all_data.shape[0] + 1, df_all_data.shape[1], left, top, width, height).table

# Populate the table with data from the DataFrame
for col_num, column_name in enumerate(df_all_data.columns):
    table.cell(0, col_num).text = column_name
    for row_num, cell_value in enumerate(df_all_data[column_name], start=1):
        table.cell(row_num, col_num).text = str(cell_value)

# Save the PowerPoint presentation to a file
prs.save("dataframe_table_example.pptx")


# In[187]:


import pandas as pd
from pptx import Presentation

from pptx.util import Inches
from pptx.dml.color import RGBColor
pptx_file_path = './June FY23 - LSP Fundamentals Q&A (English).pptx'

# Load the PowerPoint presentation
prs = Presentation(pptx_file_path)


slide = prs.slides[1]

# Get the slide layout of the current slide
slide_layout = slide.slide_layout

# Add a slide

slide = prs.slides.add_slide(slide_layout)

# Add a title
title = slide.shapes.title
title.text = "Extract QnA"

# Specify the position and dimensions of the table on the slide
left = Inches(0.2)
top = Inches(0.2)
width = Inches(16)
height = Inches(1)


# Add the DataFrame as a table to the slide
table = slide.shapes.add_table(df_all_data.shape[0] + 1, df_all_data.shape[1], left, top, width, height).table

table.columns[0].width = Inches(1.0)
table.columns[1].width = Inches(5.5)
table.columns[2].width = Inches(5.5)
table.columns[3].width = Inches(1.0)
# Populate the table with data from the DataFrame
for col_num, column_name in enumerate(df_all_data.columns):
    table.cell(0, col_num).text = column_name
    for row_num, cell_value in enumerate(df_all_data[column_name], start=1):
        table.cell(row_num, col_num).text = str(cell_value)

# Save the PowerPoint presentation to a file
prs.save("dataframe_table_example.pptx")

